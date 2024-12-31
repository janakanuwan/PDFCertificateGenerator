import pandas as pd
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from datetime import datetime
import os
from comtypes import client  # Used for PowerPoint-to-PDF conversion on Windows

# Constants
EXCEL_FILE = "Recipients.xlsx"
EXCEL_SHEET = "Details"
PPT_TEMPLATE = "Certificate_Template.pptx"
OUTPUT_DIR = "GeneratedCertificates"

# Step 1: Load the Excel File
try:
    data = pd.read_excel(EXCEL_FILE, sheet_name=EXCEL_SHEET)
    print("Excel file loaded successfully.")
except Exception as e:
    print(f"Error loading Excel file: {e}")
    exit()

# Step 2: Filter rows with non-empty ID
data = data.dropna(subset=["ID"])
if data.empty:
    print("No valid rows found in the Excel sheet.")
    exit()

# Step 3: Ensure output directory exists
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Step 4: Process each row and generate recommendation letters
for _, row in data.iterrows():
    try:
        # Extract data
        recipient_id = row["ID"]
        recipient_name = row["Name"]
        recipient_award = row["Award"]
        recipient_sponsor = row["Sponsor"]

        print(
            f"\nProcessing row ID: [{recipient_id}], Name: [{recipient_name}], Award: [{recipient_award}], Sponsor: [{recipient_sponsor}]"
        )

        # Check for NaN values and handle them explicitly
        if (pd.isna(recipient_id) or pd.isna(recipient_name) or
                not str(recipient_id).strip() or not str(recipient_name).strip()):
            print(f"Skipping row due to empty ID or Committee")
            continue

        # Convert values to strings and strip whitespace or handle NaN
        recipient_id = "" if pd.isna(recipient_id) else str(recipient_id).strip()
        recipient_name = "" if pd.isna(recipient_name) else str(recipient_name).strip()
        recipient_award = "" if pd.isna(recipient_award) else str(recipient_award).strip()
        recipient_sponsor = "" if pd.isna(recipient_sponsor) else str(recipient_sponsor).strip()

        # Open the PowerPoint template
        try:
            ppt = Presentation(PPT_TEMPLATE)
        except Exception as e:
            print(f"Error opening PowerPoint template: {e}")
            continue

        # Replace placeholders in the PowerPoint slides
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace("[[Date]]", datetime.now().strftime("%b %d, %Y"))
                            run.text = run.text.replace("[[Name]]", recipient_name)
                            run.text = run.text.replace("[[Award]]", recipient_award)
                            run.text = run.text.replace("[[Sponsor]]", recipient_sponsor)

        # Save the personalized PPT document
        ppt_file = os.path.join(OUTPUT_DIR, f"{recipient_id}.pptx")
        ppt_file = os.path.abspath(ppt_file)
        ppt.save(ppt_file)

        # Convert the PowerPoint file to PDF
        try:
            pdf_file = os.path.join(OUTPUT_DIR, f"{recipient_id}.pdf")
            pdf_file = os.path.abspath(pdf_file)
            powerpoint = client.CreateObject("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(ppt_file, WithWindow=False)
            presentation.SaveAs(pdf_file, 32)  # 32 is the format for PDF
            presentation.Close()
            powerpoint.Quit()
            print(f"Generated PDF for ID: {recipient_id}")
        except Exception as e:
            print(f"Error converting {ppt_file} to PDF: {e}")
            continue

        # Remove the intermediate PowerPoint file
        # os.remove(ppt_file)

    except KeyError as e:
        print(f"Failed to generate. Missing expected column in row: {e}")
        exit()
    except Exception as e:
        print(f"Error processing row ID {row['ID']}: {e}")

print("\nProcessing complete!")
